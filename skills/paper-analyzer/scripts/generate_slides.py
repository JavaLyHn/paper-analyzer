#!/usr/bin/env python3
"""
generate_slides.py — Academic paper slide-plan.json → .pptx

Usage:
    python3 generate_slides.py <paper-dir> [--template <template.pptx>]

paper-dir must contain:
    slide-plan.json    (written by Claude / paper-analyzer skill)
    figures/           (optional, from extract_assets.py)
    tables/            (optional)

Output:
    <paper-dir>/<output_filename>.pptx

== TEMPLATE MODE (--template <file.pptx>) ==

When the user provides a .pptx template, that file usually carries the visual
design embedded directly in its EXAMPLE SLIDES — not in the master / layouts.
Think 北工大 / corporate decks: the example slides are pre-decorated with
logos, color-block backgrounds, section-number badges, top accent bars, etc.,
which exist as ordinary `<p:sp>` / `<p:pic>` / `<p:grpSp>` elements on those
slides. The master/layouts are usually almost-blank.

Our strategy is therefore:

  1. Identify "reference slides" in the template by role:
       - cover     → typical first slide (logo + big title area)
       - section   → slide with a section-number badge (group shape)
       - content   → slide with a TITLE placeholder + top accent bar
       - closing   → the last decorative slide (often same style as section)
  2. For each item in slide-plan, DEEP-CLONE the matching reference slide
     (all its shapes/pictures/decorations are preserved), then add the new
     content (title / bullets / picture / formula) as overlays on top.
  3. After all new slides are added, delete the original reference slides.

The deep-clone is necessary because python-pptx has no public slide-copy API;
we copy the slide's `<p:sp>` tree element by element and re-establish image
and media relationships with the same rIds.

== DEFAULT MODE (no --template) ==

We fall back to the built-in navy/teal academic theme: dark-blue title bar,
white body, manual textboxes. Less polished but always available.
"""

import argparse
import copy
import json
import re
import sys
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHP
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
except ImportError:
    sys.exit("ERROR: python-pptx not installed.  Run: pip3 install python-pptx")

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None

# ── Default-theme palette (used only when no template is given) ────────────
NAVY  = RGBColor(0x1B, 0x3A, 0x6B)
TEAL  = RGBColor(0x2E, 0x86, 0xAB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x2C, 0x3E, 0x50)
GRAY  = RGBColor(0x95, 0xA5, 0xA6)
LIGHT_GRAY = RGBColor(0xBD, 0xC3, 0xC7)
LBLUE = RGBColor(0xB0, 0xC8, 0xE8)

SW_DEFAULT, SH_DEFAULT = Inches(13.33), Inches(7.5)
HDR_H_DEFAULT          = Inches(1.0)
PAD_DEFAULT            = Inches(0.4)

# ── Typography & layout discipline (applied to overlay text) ──────────────
# Use a Chinese-friendly font with reasonable English fallback. These match
# the most common academic / corporate templates in China.
FONT_LATIN = "Arial"
FONT_EA    = "微软雅黑"          # Microsoft YaHei — major font in most CN decks
FONT_EA_LIGHT = "微软雅黑 Light"  # Used for body / byline

# Type scale (in points). Keep this small — variation should come from
# weight & color, not random sizes.
SIZE_COVER_TITLE   = 40
SIZE_COVER_BYLINE  = 18
SIZE_SECTION_TITLE = 30
SIZE_BODY_TITLE    = 28
SIZE_BULLET_LARGE  = 19   # Bullets-only slide (6-8 bullets per page)
SIZE_BULLET_SMALL  = 16   # Bullets + figure (two-column, narrower lines)
SIZE_CAPTION       = 13
SIZE_FORMULA       = 30
SIZE_CLOSING       = 44   # Thank You

# Layout grid for content slides (assumes 13.33×7.5 inch / 16:9)
# These are the safe body bounds AFTER the title placeholder. The cloned
# template title placeholder ends around y=1.0 inch in the user-given deck,
# so we start body content at 1.4 inch to leave breathing room.
BODY_LEFT_PAD   = Inches(0.7)
BODY_RIGHT_PAD  = Inches(0.7)
BODY_TOP_PAD    = Inches(0.25)  # gap below title
BODY_BOTTOM_PAD = Inches(0.50)  # margin above bottom edge

# Bullet line spacing (between bullets, not within a wrapped bullet)
BULLET_SPACING_LARGE = Pt(8)
BULLET_SPACING_SMALL = Pt(5)

# Caption gap below image
CAPTION_GAP = Inches(0.10)


def _apply_font(run, *, size=None, bold=False, italic=False, color=None,
                font_latin=FONT_LATIN, font_ea=FONT_EA):
    """Apply consistent font properties to a run.

    Both the latin (西文) and east-asian (中文) typeface must be set on the
    `rPr` element — python-pptx's `font.name` only writes the latin one,
    leaving Chinese to fall through to the theme. We use lxml to set the
    `<a:ea typeface="..."/>` element so Chinese characters render in the
    intended font.
    """
    f = run.font
    if size is not None:
        f.size = Pt(size)
    f.bold = bool(bold)
    f.italic = bool(italic)
    if color is not None:
        f.color.rgb = color
    # Latin typeface
    if font_latin:
        f.name = font_latin
    # East-Asian typeface (CJK)
    if font_ea:
        rPr = run._r.get_or_add_rPr()
        # Remove any existing <a:ea>
        for ea in rPr.findall(qn("a:ea")):
            rPr.remove(ea)
        from lxml import etree
        ea_el = etree.SubElement(rPr, qn("a:ea"))
        ea_el.set("typeface", font_ea)

# Shape-element tags we copy across when cloning a slide
SHAPE_TAGS = {
    qn("p:sp"),
    qn("p:pic"),
    qn("p:grpSp"),
    qn("p:cxnSp"),
    qn("p:graphicFrame"),
}


# ────────────────────────────────────────────────────────────────────────────
# Slide cloning (deep XML copy + relationship re-bind)
# ────────────────────────────────────────────────────────────────────────────

def clone_slide(prs: Presentation, source_slide):
    """Deep-clone a slide: copy its shapes, pictures, decorations, and relationships.

    Returns the newly-added Slide.  The source slide remains untouched.

    Critical detail: rId remapping. When python-pptx creates the new slide via
    add_slide(layout), it auto-assigns `rId1` to the slide→layout relationship.
    If we then naively recreate the source slide's image relationships with
    their ORIGINAL rIds (rId1, rId2, …), we collide with the layout rel and
    PowerPoint silently fails to resolve the picture (the logo disappears).

    Solution: assign each source rel a FRESH rId in the new slide, then walk
    the cloned XML and rewrite r:embed="oldRId" → r:embed="newRId" for every
    blip / blip-extension reference.

    Implementation steps:
      1. Add a blank slide using the source's layout (this consumes one rId,
         typically rId1, for the layout).
      2. Strip the layout-inherited default shapes from the new slide.
      3. Build a {old_rId: new_rId} map by adding each source rel to the new
         slide's rels collection (let python-pptx pick the next free rId).
      4. Deep-copy source shape XML, then rewrite r:embed / r:link attributes
         using the map.
    """
    # 1) blank slide with same layout
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # 2) remove layout-inherited shapes from the new slide so we don't end up
    # with duplicate placeholders sitting under the cloned content
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        if child.tag in SHAPE_TAGS:
            sp_tree.remove(child)

    # 3) re-bind source relationships under fresh rIds, building an old→new map
    #
    # python-pptx ≥ 1.x exposes `_Relationships.get_or_add(reltype, target_part)`
    # (and `get_or_add_ext_rel(reltype, target_ref)` for external) which AUTO-
    # assigns the next rId and returns it. We can't choose the rId ourselves;
    # instead we record what the API hands back and build a remap table.
    src_rels = source_slide.part.rels
    new_rels = new_slide.part.rels
    rid_map = {}  # old_rId → new_rId

    for old_rId, rel in src_rels.items():
        if rel.reltype.endswith("/notesSlide"):
            continue  # don't drag notes
        try:
            if rel.is_external:
                new_rId = new_rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_rId = new_rels.get_or_add(rel.reltype, rel.target_part)
        except Exception as exc:  # noqa: BLE001
            print(f"  [warn] could not re-bind rel {old_rId} "
                  f"({rel.reltype.split('/')[-1]}): {exc}",
                  file=sys.stderr)
            continue
        if new_rId and new_rId != old_rId:
            rid_map[old_rId] = new_rId

    # 4) deep-copy source shape XML, rewriting rIds where they appear
    src_sp_tree = source_slide.shapes._spTree
    R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    embed_attr = R_NS + "embed"
    link_attr  = R_NS + "link"
    id_attr    = R_NS + "id"

    def _remap_rids(el):
        """Walk an XML element and rewrite r:embed / r:link / r:id values."""
        for attr in (embed_attr, link_attr, id_attr):
            v = el.get(attr)
            if v in rid_map:
                el.set(attr, rid_map[v])
        for child in el:
            _remap_rids(child)

    for el in list(src_sp_tree):
        if el.tag not in SHAPE_TAGS:
            continue
        new_el = copy.deepcopy(el)
        _remap_rids(new_el)
        sp_tree.append(new_el)

    return new_slide


def delete_slide(prs: Presentation, slide_idx: int):
    """Remove a slide from the presentation by index.

    Drops the <p:sldId> entry and the relationship; the underlying slide part
    may stay in the package but PowerPoint/Keynote/LibreOffice tolerate that.
    """
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    if slide_idx >= len(sld_ids):
        return
    sld_id_el = sld_ids[slide_idx]
    rId = sld_id_el.get(qn("r:id"))
    sld_id_lst.remove(sld_id_el)
    if rId:
        try:
            prs.part.rels.pop(rId, None)
        except Exception:
            pass


# ────────────────────────────────────────────────────────────────────────────
# Reference-slide detection
# ────────────────────────────────────────────────────────────────────────────

def _shape_has_picture(slide) -> bool:
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    return False


def _shape_has_group(slide) -> bool:
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            return True
    return False


def _has_title_placeholder(slide) -> bool:
    return slide.shapes.title is not None


def identify_references(prs: Presentation) -> dict:
    """Classify template slides into roles: cover / section / content / closing.

    Heuristic:
      - cover  : slide with embedded PICTURE shapes (logos)
                 AND no GROUP-style section badge.
                 Fallback: the first slide.
      - content: slide with a TITLE placeholder and **no** group/section badge.
                 Fallback: any slide with a TITLE placeholder.
      - section: slide that contains a GROUP shape (the "01" / "02" badge).
                 Fallback: any non-cover, non-content slide.
      - closing: the LAST section-style slide (so 01/02-numbered templates
                 reserve the highest-numbered one as a "thank you" cover).
                 Fallback: same as section.

    Returns a dict mapping role → slide index (0-based). Missing roles fall
    back to the closest available option, never failing.
    """
    n = len(prs.slides)
    if n == 0:
        return {}

    cover_idx = None
    content_idx = None
    section_idxs = []

    for i, s in enumerate(prs.slides):
        has_pic = _shape_has_picture(s)
        has_group = _shape_has_group(s)
        has_title = _has_title_placeholder(s)

        if has_title and not has_group:
            if content_idx is None:
                content_idx = i
            # later content-style slides also OK; we keep the first one
            continue
        if has_group:
            section_idxs.append(i)
            continue
        if has_pic and cover_idx is None:
            cover_idx = i
            continue

    # Fallbacks
    if cover_idx is None:
        cover_idx = 0
    if content_idx is None:
        # First slide that has a title placeholder
        for i, s in enumerate(prs.slides):
            if _has_title_placeholder(s):
                content_idx = i
                break
        if content_idx is None:
            content_idx = cover_idx
    if not section_idxs:
        # No badge-style slide — use cover as section divider too
        section_idxs = [cover_idx]

    closing_idx = section_idxs[-1]

    return {
        "cover":   cover_idx,
        "section": section_idxs[0],
        "content": content_idx,
        "closing": closing_idx,
    }


# ────────────────────────────────────────────────────────────────────────────
# Overlay helpers: add textboxes / pictures on top of a cloned slide
# ────────────────────────────────────────────────────────────────────────────

def _add_textbox(slide, left, top, width, height, text="",
                 size=18, bold=False, italic=False,
                 color=None, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP,
                 font_ea=FONT_EA, font_latin=FONT_LATIN,
                 margin=Inches(0.05)):
    """Add a textbox with consistent CJK-aware typography.

    - Sets BOTH latin and east-asian typeface so Chinese characters render in
      the intended font (微软雅黑 by default).
    - Word-wrap is on by default.
    - `anchor` controls vertical alignment (TOP / MIDDLE / BOTTOM).
    - `margin` reduces the inset around the text frame so the visible text
      lines up tightly with `(left, top, width, height)`.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = margin
    tf.margin_bottom = margin
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _apply_font(r, size=size, bold=bold, italic=italic, color=color,
                font_latin=font_latin, font_ea=font_ea)
    return tb


def _add_bullets(slide, left, top, width, height, items,
                 size=20, color=None, line_space=None,
                 anchor=MSO_ANCHOR.TOP,
                 font_ea=FONT_EA, font_latin=FONT_LATIN):
    """Add a bullet list with even line spacing & consistent typography.

    Items are separated by `line_space` (default scales with size). The
    bullet glyph is a U+2022 dot with a non-breaking space — keeps the
    bullet attached to the first word even at line breaks.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    if line_space is None:
        # Default: about 40% of the font size, in points
        line_space = Pt(size * 0.4)
    # Auto-tighten when 7+ bullets so they don't overflow the body area
    if len(items) >= 7:
        line_space = Pt(size * 0.25)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "• " + str(item)
        _apply_font(r, size=size, color=color,
                    font_latin=font_latin, font_ea=font_ea)
        p.space_after = line_space
        p.line_spacing = 1.10
    return tb


def _add_picture(slide, img_path: Path, left, top, width, height):
    """Add a picture preserving aspect ratio, centered inside the (l,t,w,h) box."""
    if PILImage is not None:
        try:
            with PILImage.open(img_path) as im:
                iw, ih = im.size
        except Exception:
            iw = ih = None
    else:
        iw = ih = None

    if iw and ih and width and height:
        img_ratio = iw / ih
        box_ratio = width / height
        if img_ratio > box_ratio:
            new_w = width
            new_h = int(width / img_ratio)
        else:
            new_h = height
            new_w = int(height * img_ratio)
        new_l = left + (width - new_w) // 2
        new_t = top + (height - new_h) // 2
        return slide.shapes.add_picture(str(img_path), new_l, new_t, new_w, new_h)
    return slide.shapes.add_picture(str(img_path), left, top, width, height)


def _set_title_placeholder(slide, text: str, *,
                           size=22, bold=True, color=None,
                           min_height=Inches(0.75),
                           font_ea=FONT_EA, font_latin=FONT_LATIN) -> bool:
    """Set the title placeholder text with disciplined typography.

    Templates often define very thin title placeholders (e.g. 0.5in tall at
    28pt). Long Chinese-English mixed titles then wrap to a second line that
    visually overflows below the title area. To avoid this we:

    1. Set an explicit, smaller font size (default 22pt) so most titles fit
       on one line at the placeholder's full width.
    2. Bump the placeholder height to `min_height` so even a 2-line title
       won't crash into the body content.
    3. Explicitly set BOTH latin and east-asian typeface for CJK fidelity.
    """
    if slide.shapes.title is None:
        return False
    title_sh = slide.shapes.title

    # Expand placeholder height if it's too cramped
    if title_sh.height is not None and title_sh.height < min_height:
        title_sh.height = min_height

    tf = title_sh.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    if p.runs:
        for r in p.runs:
            _apply_font(r, size=size, bold=bold, color=color,
                        font_latin=font_latin, font_ea=font_ea)
    return True


def _iter_text_frames(shape):
    """Recursively yield all text frames inside a shape (handles groups)."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
        if shape.shape_type == _MST.GROUP:
            for sub in shape.shapes:
                yield from _iter_text_frames(sub)
            return
    except Exception:
        pass
    if shape.has_text_frame:
        yield shape.text_frame


def _update_badge_number(slide, number):
    """Find a 1-3 digit text inside any group/textbox on the slide and replace
    it with the given number (e.g. "01" → "02"). Used for section dividers.

    Returns True if a badge text was found and updated.
    """
    target_text = f"{int(number):02d}"
    for shape in slide.shapes:
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
            if shape.shape_type == _MST.GROUP:
                for tf in _iter_text_frames(shape):
                    raw = tf.text.strip()
                    if re.fullmatch(r"\d{1,3}", raw):
                        # Replace while preserving run formatting
                        for para in tf.paragraphs:
                            for run in para.runs:
                                run.text = ""
                            if para.runs:
                                para.runs[0].text = target_text
                            else:
                                para.add_run().text = target_text
                        return True
        except Exception:
            continue
    return False


# ────────────────────────────────────────────────────────────────────────────
# Layout grid: find safe placement zones on each cloned slide
# ────────────────────────────────────────────────────────────────────────────

def _content_body_area(prs: Presentation, slide):
    """Where to put body content on a CLONED content slide.

    Body area sits below the title placeholder (or default 1.0in if missing)
    and above a bottom margin. Left/right padding is consistent with the
    title's horizontal extent when possible.
    """
    sw = prs.slide_width
    sh = prs.slide_height
    title_bottom = Inches(1.05)
    title_left   = BODY_LEFT_PAD
    title_right  = sw - BODY_RIGHT_PAD
    if slide.shapes.title is not None:
        t = slide.shapes.title
        if t.top is not None and t.height is not None:
            title_bottom = t.top + t.height
        if t.left is not None and t.width is not None:
            title_left  = t.left
            title_right = t.left + t.width

    top    = title_bottom + BODY_TOP_PAD
    bottom = sh - BODY_BOTTOM_PAD
    return (title_left, top, title_right - title_left, bottom - top)


def _cover_text_area(prs: Presentation, slide):
    """Locate the title band on a cloned cover slide.

    Strategy: find the widest, tallest AUTO_SHAPE / FREEFORM that sits in the
    middle vertical band of the slide. That's typically the brand color block.
    Returns the FULL band's (left, top, width, height) — caller decides how to
    sub-divide it for title vs byline.
    """
    sw = prs.slide_width
    sh = prs.slide_height
    best = None
    best_score = 0
    for sh_ in slide.shapes:
        if sh_.shape_type not in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
            continue
        if sh_.height is None or sh_.width is None:
            continue
        # Must be wide enough (> 60% of slide width) and not at the very top/bottom
        if sh_.width < sw * 0.6:
            continue
        cy = (sh_.top + sh_.height / 2) if sh_.top else 0
        if cy < sh * 0.15 or cy > sh * 0.85:
            continue
        score = sh_.width * sh_.height
        if score > best_score:
            best = sh_
            best_score = score
    if best is not None:
        return (best.left, best.top, best.width, best.height)
    # fallback: centered band
    return (Inches(1.0), int(sh * 0.35), sw - Inches(2.0), int(sh * 0.3))


def _section_title_anchor(prs: Presentation, slide):
    """Where to place section title text on a cloned section slide.

    Section templates usually have a GROUP shape that holds the "01" badge.
    We anchor the title BELOW that group (with breathing room), spanning the
    badge's horizontal extent so the title visually aligns with the badge.
    """
    sw = prs.slide_width
    sh = prs.slide_height
    for sh_ in slide.shapes:
        if sh_.shape_type != MSO_SHAPE_TYPE.GROUP:
            continue
        left = sh_.left
        top  = sh_.top + sh_.height + Inches(0.25)
        width = sh_.width
        return (left, top, width, Inches(1.0))
    return (Inches(1.0), int(sh * 0.55), sw - Inches(2.0), Inches(1.2))


def _find_badge_color(slide):
    """Pick a brand color from the slide's decoration shapes for accents.

    Used so section-title text matches the brand color when overlaid below
    a badge group. Returns RGBColor or None.
    """
    for sh in slide.shapes:
        if sh.shape_type != MSO_SHAPE_TYPE.GROUP:
            continue
        for sub in sh.shapes:
            try:
                if sub.fill.type == 1:  # MSO_FILL.SOLID
                    return sub.fill.fore_color.rgb
            except Exception:
                continue
    return None


# ────────────────────────────────────────────────────────────────────────────
# Slide-type renderers (template mode)
# ────────────────────────────────────────────────────────────────────────────

def _tpl_title_slide(prs, s, refs):
    """Title (cover) slide.

    Layout:
        +------------------------------------------+
        | (logo)                       (logo)       |   ← preserved from template
        |    ─────                                  |
        | ┌─────────── BLUE BAND ──────────────┐   |
        | │       PAPER TITLE  (38pt, bold)     │   |   ← upper 55% of band
        | │                                      │   |
        | │   Authors · Venue · Year (18pt)      │   |   ← lower 45% of band
        | └──────────────────────────────────────┘   |
        +------------------------------------------+
    """
    src = prs.slides[refs["cover"]]
    new_slide = clone_slide(prs, src)

    title = s.get("title", "")
    byline_parts = list(filter(None, [
        s.get("authors", ""),
        " · ".join(filter(None, [s.get("venue", ""), str(s.get("year", "")) if s.get("year") else ""])),
    ]))
    byline = "\n".join(byline_parts)

    band_left, band_top, band_w, band_h = _cover_text_area(prs, new_slide)
    inset = Inches(0.6)
    safe_left  = band_left + inset
    safe_width = band_w - inset * 2

    # Title — vertically centered in upper 60% of the band, padded from edges
    title_top    = band_top + Inches(0.35)
    title_height = int(band_h * 0.55)
    _add_textbox(new_slide, safe_left, title_top, safe_width, title_height,
                 title,
                 size=SIZE_COVER_TITLE, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Byline — middle-anchored block in lower 40%
    if byline:
        byl_top    = band_top + int(band_h * 0.62)
        byl_height = int(band_h * 0.32)
        _add_textbox(new_slide, safe_left, byl_top, safe_width, byl_height,
                     byline,
                     size=SIZE_COVER_BYLINE, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     font_ea=FONT_EA_LIGHT)
    return new_slide


def _tpl_section_slide(prs, s, refs, *, section_num):
    """Section divider.

    Layout:
        +----------------------------------+
        | (top decoration curve)            |
        |                                   |
        |          ┌────┐                   |
        |          │ 01 │                   |   ← badge group (auto-numbered)
        |          └────┘                   |
        |    Section Title (32pt, bold)     |   ← centered below badge
        |                                   |
        | (bottom decoration curve)         |
        +----------------------------------+
    """
    src = prs.slides[refs["section"]]
    new_slide = clone_slide(prs, src)
    _update_badge_number(new_slide, section_num)

    title = s.get("title", "")
    accent = _find_badge_color(new_slide) or DARK

    if not _set_title_placeholder(new_slide, title):
        left, top, width, height = _section_title_anchor(prs, new_slide)
        _add_textbox(new_slide, left, top, width, height, title,
                     size=SIZE_SECTION_TITLE, bold=True, color=accent,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    return new_slide


def _find_brand_color(prs: Presentation, refs: dict):
    """Pick the brand color from the template's section-divider badge group.

    Returns RGBColor or None. Reused for consistent title color across
    all content slides.
    """
    if "section" not in refs:
        return None
    src = prs.slides[refs["section"]]
    return _find_badge_color(src)


def _tpl_bullets_slide(prs, s, refs, paper_dir):
    """Content slide with bullets (and optional figure on the right).

    Layout (no figure):
        +-------------- TITLE ---------------+
        |  • bullet 1                         |
        |  • bullet 2                         |   ← 22pt, line-spacing 1.15,
        |  • bullet 3                         |     12pt gap between bullets
        |  • bullet 4                         |
        +-------------------------------------+

    Layout (with figure, two-column):
        +-------------- TITLE ---------------+
        |  • bullet 1         |   ┌─────┐    |
        |  • bullet 2         |   │ FIG │    |
        |  • bullet 3         |   └─────┘    |   ← 19pt bullets, image right
        +-------------------------------------+
    """
    src = prs.slides[refs["content"]]
    new_slide = clone_slide(prs, src)
    _set_title_placeholder(new_slide, s.get("title", ""))

    items   = s.get("bullets", [])
    fig_rel = s.get("figure")
    fig_path = (paper_dir / fig_rel) if fig_rel else None
    has_fig  = bool(fig_path and fig_path.exists())

    left, top, width, height = _content_body_area(prs, new_slide)

    if has_fig:
        gap = Inches(0.35)
        bw = int(width * 0.50) - int(gap / 2)
        pw = width - bw - int(gap)
        _add_bullets(new_slide, left, top, bw, height, items,
                     size=SIZE_BULLET_SMALL, color=DARK,
                     line_space=BULLET_SPACING_SMALL,
                     anchor=MSO_ANCHOR.MIDDLE)
        _add_picture(new_slide, fig_path,
                     left + bw + gap, top, pw, height)
    else:
        _add_bullets(new_slide, left, top, width, height, items,
                     size=SIZE_BULLET_LARGE, color=DARK,
                     line_space=BULLET_SPACING_LARGE,
                     anchor=MSO_ANCHOR.TOP)
    return new_slide


def _tpl_image_slide(prs, s, refs, paper_dir):
    """Image-only content slide with caption.

    Layout:
        +-------------- TITLE ---------------+
        |          ┌─────────────┐            |
        |          │             │            |
        |          │   FIGURE    │            |   ← max body height − caption
        |          │             │            |
        |          └─────────────┘            |
        |   Figure N — caption (14pt italic)  |   ← 0.5in bottom strip
        +-------------------------------------+
    """
    src = prs.slides[refs["content"]]
    new_slide = clone_slide(prs, src)
    _set_title_placeholder(new_slide, s.get("title", ""))

    img_rel  = s.get("image")
    img_path = (paper_dir / img_rel) if img_rel else None
    caption  = s.get("caption", "")

    left, top, width, height = _content_body_area(prs, new_slide)
    cap_h = Inches(0.5) if caption else Inches(0)
    pic_h = height - cap_h - (CAPTION_GAP if caption else Inches(0))

    if img_path and img_path.exists():
        _add_picture(new_slide, img_path, left, top, width, pic_h)
    else:
        _add_textbox(new_slide, left, top, width, pic_h,
                     f"[图像未找到 / Image not found: {img_rel}]",
                     size=16, color=GRAY, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
    if caption:
        _add_textbox(new_slide, left, top + pic_h + CAPTION_GAP,
                     width, cap_h, caption,
                     size=SIZE_CAPTION, italic=True, color=GRAY,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
                     font_ea=FONT_EA_LIGHT)
    return new_slide


def _render_latex(latex: str, color="#2C3E50"):
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        expr = latex.strip()
        if not expr.startswith("$"):
            expr = f"${expr}$"
        fig = plt.figure(figsize=(11, 1.8), facecolor="none")
        fig.text(0.5, 0.5, expr, ha="center", va="center",
                 fontsize=26, color=color)
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        fig.savefig(tmp.name, bbox_inches="tight", transparent=True, dpi=150)
        plt.close(fig)
        return tmp.name
    except Exception as e:
        print(f"  [warn] LaTeX render failed: {e}", file=sys.stderr)
        return None


def _tpl_formula_slide(prs, s, refs, paper_dir):
    """Formula slide: centered LaTeX render + caption below.

    Layout:
        +-------------- TITLE ---------------+
        |                                     |
        |                                     |
        |          F O R M U L A              |   ← rendered PNG, vert-center
        |                                     |
        |                                     |
        |     caption explaining formula      |   ← 17pt, centered
        +-------------------------------------+
    """
    src = prs.slides[refs["content"]]
    new_slide = clone_slide(prs, src)
    _set_title_placeholder(new_slide, s.get("title", ""))

    latex   = s.get("latex", "")
    caption = s.get("caption", "")

    left, top, width, height = _content_body_area(prs, new_slide)
    cap_h = Inches(0.8) if caption else Inches(0)
    # Reserve roughly 2/3 of remaining height for the formula image
    formula_h = int((height - cap_h) * 0.7)
    # Vertically center the formula in (top, top + height - cap_h)
    formula_block = height - cap_h
    formula_top = top + (formula_block - formula_h) // 2

    formula_img = _render_latex(latex)
    if formula_img:
        _add_picture(new_slide, formula_img, left, formula_top, width, formula_h)
    else:
        _add_textbox(new_slide, left, formula_top, width, formula_h, latex,
                     size=SIZE_FORMULA, color=DARK,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if caption:
        _add_textbox(new_slide, left, top + formula_block + Inches(0.15),
                     width, cap_h - Inches(0.15), caption,
                     size=17, color=DARK,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
                     font_ea=FONT_EA_LIGHT)
    return new_slide


def _tpl_closing_slide(prs, s, refs, *, section_num):
    """Closing (thank-you) slide. Reuses the section-divider template.

    Layout: same as section, with "Thank You" replacing the section title and
    an optional subtitle below.
    """
    src = prs.slides[refs["closing"]]
    new_slide = clone_slide(prs, src)
    _update_badge_number(new_slide, section_num)

    title = s.get("title", "Thank You")
    accent = _find_badge_color(new_slide) or DARK

    if not _set_title_placeholder(new_slide, title):
        left, top, width, height = _section_title_anchor(prs, new_slide)
        _add_textbox(new_slide, left, top, width, Inches(1.3), title,
                     size=SIZE_CLOSING, bold=True, color=accent,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        if s.get("subtitle"):
            _add_textbox(new_slide, left, top + Inches(1.3),
                         width, Inches(0.6), s["subtitle"],
                         size=18, color=GRAY,
                         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                         font_ea=FONT_EA_LIGHT)
    return new_slide


# ────────────────────────────────────────────────────────────────────────────
# Default-theme renderers (no --template)
# ────────────────────────────────────────────────────────────────────────────

def _df_solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _df_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _df_header(slide, title, sw):
    bar = slide.shapes.add_shape(SHP.RECTANGLE, 0, 0, sw, HDR_H_DEFAULT)
    _df_solid(bar, NAVY)
    _add_textbox(slide, PAD_DEFAULT, Inches(0.1),
                 sw - PAD_DEFAULT * 2, HDR_H_DEFAULT - Inches(0.1),
                 title, size=28, bold=True, color=WHITE)


def _df_title_slide(prs, s):
    sw, sh = prs.slide_width, prs.slide_height
    slide = _df_blank(prs)
    bg = slide.shapes.add_shape(SHP.RECTANGLE, 0, 0, sw, sh); _df_solid(bg, NAVY)
    _add_textbox(slide, PAD_DEFAULT, Inches(1.6),
                 sw - PAD_DEFAULT * 2, Inches(2.4),
                 s.get("title", ""), size=32, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
    byline = "\n".join(filter(None, [
        s.get("authors", ""),
        " · ".join(filter(None, [s.get("venue", ""), str(s.get("year", "")) if s.get("year") else ""])),
    ]))
    _add_textbox(slide, PAD_DEFAULT, Inches(4.3),
                 sw - PAD_DEFAULT * 2, Inches(1.4),
                 byline, size=18, color=LBLUE, align=PP_ALIGN.CENTER)


def _df_section_slide(prs, s):
    sw, sh = prs.slide_width, prs.slide_height
    slide = _df_blank(prs)
    accent = slide.shapes.add_shape(SHP.RECTANGLE, 0, 0, Inches(4.5), sh)
    _df_solid(accent, NAVY)
    _add_textbox(slide, Inches(5.0), Inches(2.7),
                 sw - Inches(5.4), Inches(2.1),
                 s.get("title", ""), size=36, bold=True, color=DARK)


def _df_bullets_slide(prs, s, paper_dir):
    sw, sh = prs.slide_width, prs.slide_height
    slide = _df_blank(prs); _df_header(slide, s.get("title", ""), sw)
    items = s.get("bullets", [])
    fig_rel = s.get("figure"); fig_path = (paper_dir / fig_rel) if fig_rel else None
    body_y = HDR_H_DEFAULT + PAD_DEFAULT
    body_h = sh - body_y - PAD_DEFAULT
    body_w = sw - PAD_DEFAULT * 2
    if fig_path and fig_path.exists():
        _add_bullets(slide, PAD_DEFAULT, body_y, Inches(7.0), body_h, items,
                     color=DARK, size=19)
        _add_picture(slide, fig_path, Inches(7.6), body_y, Inches(5.3), body_h)
    else:
        _add_bullets(slide, PAD_DEFAULT, body_y, body_w, body_h, items,
                     color=DARK, size=20)


def _df_image_slide(prs, s, paper_dir):
    sw, sh = prs.slide_width, prs.slide_height
    slide = _df_blank(prs); _df_header(slide, s.get("title", ""), sw)
    caption = s.get("caption", "")
    body_y = HDR_H_DEFAULT + PAD_DEFAULT
    body_w = sw - PAD_DEFAULT * 2
    cap_h = Inches(0.42) if caption else Inches(0)
    body_h = sh - body_y - PAD_DEFAULT - cap_h
    img_rel = s.get("image"); img_path = (paper_dir / img_rel) if img_rel else None
    if img_path and img_path.exists():
        _add_picture(slide, img_path, PAD_DEFAULT, body_y, body_w, body_h)
    else:
        _add_textbox(slide, PAD_DEFAULT, body_y, body_w, body_h,
                     f"[图像未找到: {img_rel}]", size=16, color=GRAY)
    if caption:
        _add_textbox(slide, PAD_DEFAULT, body_y + body_h + Inches(0.05),
                     body_w, cap_h, caption,
                     size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


def _df_formula_slide(prs, s, paper_dir):
    sw, sh = prs.slide_width, prs.slide_height
    slide = _df_blank(prs); _df_header(slide, s.get("title", ""), sw)
    latex = s.get("latex", ""); caption = s.get("caption", "")
    body_w = sw - PAD_DEFAULT * 2
    y = HDR_H_DEFAULT + PAD_DEFAULT
    formula_img = _render_latex(latex)
    if formula_img:
        _add_picture(slide, formula_img, PAD_DEFAULT, y, body_w, Inches(2.4))
        y += Inches(2.7)
    else:
        _add_textbox(slide, PAD_DEFAULT, y, body_w, Inches(1.6),
                     latex, size=20, color=DARK, align=PP_ALIGN.CENTER)
        y += Inches(1.9)
    if caption:
        _add_textbox(slide, PAD_DEFAULT, y, body_w, sh - y - PAD_DEFAULT,
                     caption, size=17, color=DARK)


def _df_closing_slide(prs, s):
    sw, sh = prs.slide_width, prs.slide_height
    slide = _df_blank(prs)
    bg = slide.shapes.add_shape(SHP.RECTANGLE, 0, 0, sw, sh); _df_solid(bg, NAVY)
    _add_textbox(slide, PAD_DEFAULT, Inches(2.4),
                 sw - PAD_DEFAULT * 2, Inches(1.8),
                 s.get("title", "Thank You"),
                 size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if s.get("subtitle"):
        _add_textbox(slide, PAD_DEFAULT, Inches(4.4),
                     sw - PAD_DEFAULT * 2, Inches(1.2),
                     s["subtitle"], size=26, color=LBLUE, align=PP_ALIGN.CENTER)


# ────────────────────────────────────────────────────────────────────────────
# Main
# ────────────────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("paper_dir", help="Directory with slide-plan.json and figures/")
    ap.add_argument("--template", default=None,
                    help="User-provided .pptx template (clones its example slides)")
    args = ap.parse_args()

    paper_dir = Path(args.paper_dir).resolve()
    plan_path = paper_dir / "slide-plan.json"
    if not plan_path.exists():
        sys.exit(f"ERROR: slide-plan.json not found in {paper_dir}")

    with open(plan_path, encoding="utf-8") as f:
        plan = json.load(f)

    use_template = bool(args.template)
    slides = plan.get("slides", [])

    if use_template:
        tpl = Path(args.template).resolve()
        if not tpl.exists():
            sys.exit(f"ERROR: Template file not found: {tpl}")
        prs = Presentation(str(tpl))
        n_template_slides = len(prs.slides)
        refs = identify_references(prs)
        print(f"Using template: {tpl.name}")
        print(f"  template has {n_template_slides} reference slides")
        print(f"  reference roles: cover=slide{refs['cover']+1}, "
              f"section=slide{refs['section']+1}, "
              f"content=slide{refs['content']+1}, "
              f"closing=slide{refs['closing']+1}")
        print(f"Generating {len(slides)} new slides by cloning reference slides …")

        section_counter = 0
        for i, s in enumerate(slides):
            stype = s.get("type", "bullets")
            if stype == "title":
                _tpl_title_slide(prs, s, refs)
            elif stype == "section":
                section_counter += 1
                _tpl_section_slide(prs, s, refs, section_num=section_counter)
            elif stype == "bullets":
                _tpl_bullets_slide(prs, s, refs, paper_dir)
            elif stype == "image":
                _tpl_image_slide(prs, s, refs, paper_dir)
            elif stype == "formula":
                _tpl_formula_slide(prs, s, refs, paper_dir)
            elif stype == "closing":
                _tpl_closing_slide(prs, s, refs, section_num=section_counter + 1)
            else:
                print(f"  [warn] Unknown slide type '{stype}' — skipped", file=sys.stderr)
                continue
            print(f"  [{i+1}/{len(slides)}] {stype}: {s.get('title', '')[:60]}")

        # Now delete the ORIGINAL reference slides (the first N_template_slides slides)
        # Delete from the END so indexes stay valid.
        for orig_idx in range(n_template_slides - 1, -1, -1):
            delete_slide(prs, orig_idx)
    else:
        prs = Presentation()
        prs.slide_width  = SW_DEFAULT
        prs.slide_height = SH_DEFAULT
        print("Using default academic theme (navy / teal)")
        print(f"Generating {len(slides)} slides …")
        for i, s in enumerate(slides):
            stype = s.get("type", "bullets")
            if stype == "title":
                _df_title_slide(prs, s)
            elif stype == "section":
                _df_section_slide(prs, s)
            elif stype == "bullets":
                _df_bullets_slide(prs, s, paper_dir)
            elif stype == "image":
                _df_image_slide(prs, s, paper_dir)
            elif stype == "formula":
                _df_formula_slide(prs, s, paper_dir)
            elif stype == "closing":
                _df_closing_slide(prs, s)
            else:
                print(f"  [warn] Unknown slide type '{stype}' — skipped", file=sys.stderr)
                continue
            print(f"  [{i+1}/{len(slides)}] {stype}: {s.get('title', '')[:60]}")

    out_stem = plan.get("output_filename", "slides")
    out_path = paper_dir / (out_stem + ".pptx")
    prs.save(str(out_path))
    print(f"\n✓ Saved: {out_path}  ({len(slides)} slides)")


if __name__ == "__main__":
    main()
